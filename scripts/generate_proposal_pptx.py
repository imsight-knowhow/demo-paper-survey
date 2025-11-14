from pptx import Presentation


def add_title_and_bullets(presentation: Presentation, title: str, bullets: list[str]) -> None:
    slide_layout = presentation.slide_layouts[1]  # title and content
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    text_frame = body.text_frame
    text_frame.clear()

    for index, bullet in enumerate(bullets):
        if index == 0 and text_frame.paragraphs:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0


def main() -> None:
    presentation = Presentation()

    # 1. 封面
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "StreamDiffusionV2 实时视频扩散平台"
    subtitle = slide.placeholders[1]
    subtitle.text = "商业计划概览\nJames · 创业导师视角"

    # 2. 市场与机会
    add_title_and_bullets(
        presentation,
        "市场与机会",
        [
            "全球直播与互动视频市场高速增长，规模达数百亿美元",
            "虚拟主播、互动综艺、直播电商对视频智能需求强烈",
            "视频扩散模型从离线内容制作走向实时生成",
        ],
    )

    # 3. 核心痛点
    add_title_and_bullets(
        presentation,
        "核心痛点",
        [
            "离线视频扩散模型难以直接满足直播级延迟与稳定性",
            "现有虚拟人方案多依赖传统 2D/3D 与动捕，缺乏生成多样性",
            "平台自研需要同时掌握模型、GPU 并行与流媒体工程，成本极高",
        ],
    )

    # 4. 解决方案概览
    add_title_and_bullets(
        presentation,
        "解决方案概览",
        [
            "打造面向直播与互动视频的实时视频扩散基础设施平台",
            "在不改模型权重的前提下，通过系统级优化实现低延迟与高稳定性",
            "一体化接入推流/拉流协议、SLO 调度与监控告警",
        ],
    )

    # 5. 核心技术优势
    add_title_and_bullets(
        presentation,
        "核心技术优势",
        [
            "SLO-aware Batching Scheduler：在延迟预算内最大化 GPU 利用率",
            "Pipeline-parallel Stream-Batch：多 GPU 流水线并行，接近线性加速",
            "Rolling KV cache + sink tokens + RoPE 刷新：支持长时稳定直播",
            "Motion-aware noise schedule：在画质与实时性之间自适应平衡",
        ],
    )

    # 6. 产品形态
    add_title_and_bullets(
        presentation,
        "产品形态",
        [
            "云端实时推理服务：API/流式接口，按使用量计费",
            "企业版引擎：容器化部署，支持自定义模型与私有化集群",
            "创作者与主播工具：Studio/插件/OBS 集成，一键开启风格化直播",
        ],
    )

    # 7. 商业模式
    add_title_and_bullets(
        presentation,
        "商业模式",
        [
            "SaaS/PaaS：按 GPU 计算量、帧数与分辨率计费，提供套餐与 SLA",
            "企业 License：按节点数/并发流数收取年费，含技术支持与升级",
            "项目制：虚拟演唱会、虚拟展览等高客单价定制项目",
        ],
    )

    # 8. 典型客户与场景
    add_title_and_bullets(
        presentation,
        "典型客户与场景",
        [
            "虚拟主播机构与内容工厂：统一风格的虚拟形象与场景切换",
            "直播电商品牌方：品牌滤镜与数字导购员，提升辨识度与转化率",
            "虚拟演唱会与互动活动：高密度特效与沉浸式虚拟场景",
        ],
    )

    # 9. 竞争与差异化
    add_title_and_bullets(
        presentation,
        "竞争与差异化",
        [
            "区别于 Runway/Pika 等离线视频平台，专注实时直播场景",
            "区别于传统虚拟人方案，以生成式视频扩散提供更高表达力",
            "系统工程 + 调度一体化设计，可插拔模型，易与平台/云厂商合作",
        ],
    )

    # 10. 里程碑与风险应对
    add_title_and_bullets(
        presentation,
        "里程碑与风险应对",
        [
            "12–18 个月：从单路 Demo 到多路并发 MVP，再到商用 PoC 与规模化",
            "技术与成本风险：持续优化调度与模型，分层产品档位控制毛利",
            "法律与市场风险：提前做专利 FTO 检索，先深耕虚拟主播与电商场景",
        ],
    )

    presentation.save("proposal.pptx")


if __name__ == "__main__":
    main()

